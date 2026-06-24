"""
Lab Meeting Progress Report — SeqSQLi
Run:
    python3 presentasi/build_ppt_progress.py
Output: presentasi/SeqSQLi_ProgressReport.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT   = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIGDIR = os.path.join(ROOT, "figures")
OUT    = os.path.join(ROOT, "presentasi", "SeqSQLi_ProgressReport.pptx")

BLUE_D  = RGBColor(0x21, 0x66, 0xac)
BLUE_L  = RGBColor(0xdc, 0xe8, 0xf5)
ORANGE  = RGBColor(0xe0, 0x8d, 0x00)
RED     = RGBColor(0xd6, 0x60, 0x4d)
GREEN   = RGBColor(0x2e, 0x8b, 0x57)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x1a, 0x1a, 0x1a)
GRAY    = RGBColor(0x55, 0x55, 0x55)
LGRAY   = RGBColor(0xcc, 0xcc, 0xcc)
BGWHITE = RGBColor(0xF7, 0xF9, 0xFC)
YELLOW  = RGBColor(0xff, 0xf3, 0xd6)

SW = Inches(13.33)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def bg(slide, color=BGWHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def txbox(slide, text, x, y, w, h, size=18, bold=False,
          color=BLACK, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def rect(slide, x, y, w, h, fill, ec=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if ec:
        s.line.color.rgb = ec
        s.line.width = lw
    else:
        s.line.fill.background()
    return s


def hbar(slide, title, subtitle=None):
    rect(slide, 0, 0, SW, Inches(1.3), BLUE_D)
    txbox(slide, title, Inches(0.45), Inches(0.1),
          Inches(12.4), Inches(0.75), size=30, bold=True, color=WHITE)
    if subtitle:
        txbox(slide, subtitle, Inches(0.45), Inches(0.8),
              Inches(12.4), Inches(0.42), size=15, color=BLUE_L)


def pic(slide, path, x, y, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, w, h)
    else:
        txbox(slide, "[missing: " + os.path.basename(path) + "]",
              x, y, w, h, size=12, color=RED)


def status_badge(slide, x, y, label, done=True):
    color = GREEN if done else ORANGE
    rect(slide, x, y, Inches(1.5), Inches(0.38), color)
    txbox(slide, label, x, y + Inches(0.04), Inches(1.5), Inches(0.35),
          size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ─── SLIDE 1: TITLE ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, BLUE_D)
rect(s, 0, Inches(2.3), SW, Inches(3.4), RGBColor(0x19, 0x55, 0x9a))
txbox(s, "SeqSQLi", Inches(1), Inches(0.85), Inches(11.3), Inches(1.1),
      size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, "Sequential SQL Injection WAF Bypass via Deep Reinforcement Learning",
      Inches(1), Inches(1.95), Inches(11.3), Inches(0.7),
      size=20, color=BLUE_L, align=PP_ALIGN.CENTER)
rect(s, Inches(3.5), Inches(2.95), Inches(6.3), Inches(0.05),
     RGBColor(0x88, 0xaa, 0xcc))
txbox(s, "PROGRESS REPORT  —  Lab Meeting",
      Inches(1), Inches(3.1), Inches(11.3), Inches(0.65),
      size=22, bold=True, color=RGBColor(0xff, 0xe0, 0x80), align=PP_ALIGN.CENTER)
txbox(s, "Roby Firnando Yusuf  |  Advisor: Prof. Duksan Ryu  |  June 2026",
      Inches(1), Inches(4.0), Inches(11.3), Inches(0.55),
      size=16, color=RGBColor(0xaa, 0xcc, 0xee), align=PP_ALIGN.CENTER)


# ─── SLIDE 2: AGENDA ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
hbar(s, "Today's Agenda")
items = [
    ("1.  Research Overview", BLACK),
    ("2.  RQ1 — Algorithm Comparison Results  ✅", GREEN),
    ("3.  RQ3 — Mutation Ordering Analysis  ✅", GREEN),
    ("4.  Paper Writing Status", BLUE_D),
    ("5.  Next Steps & Timeline", ORANGE),
]
for i, (txt, col) in enumerate(items):
    txbox(s, txt, Inches(2.5), Inches(1.7 + i * 0.92), Inches(9), Inches(0.75),
          size=24, bold=True, color=col)


# ─── SLIDE 3: RESEARCH OVERVIEW ──────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
hbar(s, "Research Overview", "SeqSQLi: RL-based sequential mutation for WAF bypass")

txbox(s, "Problem", Inches(0.4), Inches(1.45), Inches(4.0), Inches(0.45),
      size=15, bold=True, color=RED)
txbox(s, "WAFs can be evaded by changing SQL payload syntax while preserving semantics.\n"
         "Existing approaches do not learn from live WAF feedback or study mutation ordering.",
      Inches(0.4), Inches(1.9), Inches(4.1), Inches(1.5), size=14, color=BLACK)

txbox(s, "Approach", Inches(4.7), Inches(1.45), Inches(4.0), Inches(0.45),
      size=15, bold=True, color=BLUE_D)
txbox(s, "RL agent selects from 51 mutation operators per step.\n"
         "Live WAF response (200/403) = reward signal.\n"
         "Episode: 1 payload, max 15 steps.",
      Inches(4.7), Inches(1.9), Inches(4.1), Inches(1.5), size=14, color=BLACK)

txbox(s, "Algorithms", Inches(9.1), Inches(1.45), Inches(3.9), Inches(0.45),
      size=15, bold=True, color=ORANGE)
txbox(s, "PPO  ·  TRPO  ·  A2C\nAll on-policy, identical setup.\n150k timesteps each.",
      Inches(9.1), Inches(1.9), Inches(3.9), Inches(1.5), size=14, color=BLACK)

rect(s, Inches(0.35), Inches(3.6), Inches(12.6), Inches(0.04), LGRAY)

txbox(s, "Research Questions", Inches(0.4), Inches(3.75), Inches(12.5), Inches(0.45),
      size=15, bold=True, color=BLUE_D)
rqs = [
    ("RQ1", "Which deep-RL algorithm is most effective & sample-efficient?", "✅ DONE"),
    ("RQ2", "Can the policy transfer from ModSecurity → Safeline CE (ML-based WAF)?", "⏳ NEXT"),
    ("RQ3", "Does mutation ordering significantly affect bypass success rate?", "✅ DONE"),
]
for i, (rq, q, status) in enumerate(rqs):
    col = GREEN if "DONE" in status else ORANGE
    rect(s, Inches(0.35), Inches(4.3 + i * 0.88), Inches(12.6), Inches(0.78),
         BGWHITE, col, Pt(1))
    txbox(s, rq, Inches(0.5), Inches(4.37 + i * 0.88), Inches(0.7), Inches(0.6),
          size=14, bold=True, color=col)
    txbox(s, q, Inches(1.3), Inches(4.37 + i * 0.88), Inches(10.2), Inches(0.6),
          size=14, color=BLACK)
    txbox(s, status, Inches(11.6), Inches(4.37 + i * 0.88), Inches(1.2), Inches(0.6),
          size=13, bold=True, color=col, align=PP_ALIGN.RIGHT)


# ─── SLIDE 4: RQ1 RESULTS ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
hbar(s, "RQ1 — Algorithm Comparison Results  ✅",
     "Union corpus: 108 payloads × 3 complexity tiers  |  150k timesteps  |  ModSecurity CRS v3.3.2")

col_hdrs = ["Algorithm", "IFNR ↑", "SPBARC ↓", "Trivial", "Medium", "Complex"]
col_ws   = [Inches(2.3), Inches(1.85), Inches(1.95), Inches(1.55), Inches(1.55), Inches(1.7)]
x0, y0  = Inches(0.5), Inches(1.48)

x = x0
for hdr, w in zip(col_hdrs, col_ws):
    rect(s, x, y0, w, Inches(0.58), BLUE_D)
    txbox(s, hdr, x + Inches(0.05), y0 + Inches(0.06), w - Inches(0.1), Inches(0.46),
          size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    x += w

rows = [
    ("TRPO", "99.1%", "6.07",  "100%", "100%", "97.2%", BLUE_D, BLUE_L),
    ("PPO",  "88.9%", "6.89",  "100%", "100%", "66.7%", ORANGE, YELLOW),
    ("A2C",  "76.9%", "10.08", "100%", "97.2%", "33.3%", RED,   RGBColor(0xfd, 0xf0, 0xee)),
    ("Random (baseline)", "3.7%", "216", "—", "—", "—", GRAY, RGBColor(0xf0, 0xf0, 0xf0)),
]
for i, (algo, ifnr, spbarc, triv, med, comp, ec, fc) in enumerate(rows):
    x = x0
    vals = [algo, ifnr, spbarc, triv, med, comp]
    for j, (val, w) in enumerate(zip(vals, col_ws)):
        rect(s, x, y0 + Inches(0.58 + i * 0.82), w, Inches(0.78), fc)
        bold = j in [0, 5] and i < 3
        txbox(s, val, x + Inches(0.05), y0 + Inches(0.66 + i * 0.82),
              w - Inches(0.1), Inches(0.62), size=17, bold=bold,
              color=ec if j in [0, 5] else BLACK, align=PP_ALIGN.CENTER)
        x += w

txbox(s, "IFNR = Improved False Negative Rate  (% of initially blocked payloads bypassed)\n"
         "SPBARC = Successful Payload Bypass Avg Request Count  (lower = more efficient)",
      Inches(0.5), Inches(5.05), Inches(12.3), Inches(0.55), size=11, color=GRAY, italic=True)

rect(s, Inches(0.35), Inches(5.7), Inches(12.6), Inches(1.42), BLUE_L, BLUE_D, Pt(1))
txbox(s, "Key Finding: Complex tier is the discriminator",
      Inches(0.55), Inches(5.76), Inches(12.1), Inches(0.5), size=15, bold=True, color=BLUE_D)
txbox(s, "All 3 algorithms reach ~100% on Trivial & Medium. "
         "Performance diverges sharply on Complex (97.2% vs 66.7% vs 33.3%).\n"
         "Complex payloads require longer, order-constrained mutation chains — "
         "where TRPO's hard KL-bound constraint provides the critical stability advantage.",
      Inches(0.55), Inches(6.24), Inches(12.1), Inches(0.8), size=13, color=BLACK)


# ─── SLIDE 5: TIER CHART ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
hbar(s, "RQ1 — Per-Tier Breakdown",
     "Why trivial and medium are saturated; complex is the real challenge")
pic(s, os.path.join(FIGDIR, "fig_tier_barchart.png"),
    Inches(0.3), Inches(1.38), Inches(8.2), Inches(5.7))
txbox(s, "Tier Definition", Inches(8.7), Inches(1.5), Inches(4.3), Inches(0.45),
      size=15, bold=True, color=BLUE_D)
tier_defs = [
    ("Trivial", "Case + whitespace only", BLACK),
    ("Medium",  "+ Aggregate fn rule\n(GROUP_CONCAT blocked)", BLACK),
    ("Complex", "+ Hex encoding + FROM users\n→ 3 rule families simultaneously", RED),
]
for i, (tier, desc, col) in enumerate(tier_defs):
    rect(s, Inches(8.7), Inches(2.05 + i * 1.5), Inches(4.3), Inches(1.35),
         BGWHITE, BLUE_D if col == BLACK else RED, Pt(1))
    txbox(s, tier, Inches(8.85), Inches(2.12 + i * 1.5), Inches(4.0), Inches(0.48),
          size=14, bold=True, color=BLUE_D if col == BLACK else RED)
    txbox(s, desc, Inches(8.85), Inches(2.58 + i * 1.5), Inches(4.0), Inches(0.75),
          size=13, color=BLACK)
txbox(s, "agg_swap mutation (GROUP_CONCAT→JSON_ARRAYAGG)\nunlocked medium tier — was 0% before",
      Inches(8.7), Inches(6.65), Inches(4.3), Inches(0.72),
      size=12, color=GREEN, italic=True, bold=True)


# ─── SLIDE 6: RQ3 ORDERING ───────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
hbar(s, "RQ3 — Does Mutation Order Matter?  ✅",
     "Answer: Yes — ordering encodes structural signal from WAF rule topology")
pic(s, os.path.join(FIGDIR, "fig_ordering_pairs.png"),
    Inches(0.3), Inches(1.38), Inches(8.4), Inches(5.8))
txbox(s, "Summary", Inches(8.95), Inches(1.55), Inches(4.05), Inches(0.45),
      size=15, bold=True, color=BLUE_D)
findings = [
    ("35 consensus pairs\n(≥2 algorithms agree)", BLUE_D),
    ("Strongest:\nident_backtick→hex_to_char\nForward: 98%  |  Reversed: 0%", RED),
    ("Pattern consistent\nacross PPO, TRPO, A2C", GREEN),
    ("Conclusion: causal from\nWAF rule structure,\nnot learning artifact", BLUE_D),
]
for i, (txt, col) in enumerate(findings):
    rect(s, Inches(8.95), Inches(2.1 + i * 1.28), Inches(4.1), Inches(1.18),
         BGWHITE, col, Pt(1.2))
    txbox(s, txt, Inches(9.1), Inches(2.17 + i * 1.28), Inches(3.8), Inches(1.05),
          size=13, color=BLACK)


# ─── SLIDE 7: PAPER STATUS ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
hbar(s, "Paper Writing Status", "Two targets: Final Exam & Publication (KIIT Journal)")

# Final Exam column
rect(s, Inches(0.35), Inches(1.4), Inches(6.1), Inches(5.72), BGWHITE, BLUE_D, Pt(1.5))
txbox(s, "Final Exam Paper", Inches(0.55), Inches(1.5), Inches(5.7), Inches(0.5),
      size=18, bold=True, color=BLUE_D)
txbox(s, "docs/manuscript-final-exam.hwp", Inches(0.55), Inches(1.98),
      Inches(5.7), Inches(0.38), size=11, color=GRAY, italic=True)
txbox(s, "Deadline: ~16 June 2026", Inches(0.55), Inches(2.34),
      Inches(5.7), Inches(0.38), size=13, bold=True, color=ORANGE)

final_items = [
    ("All 5 sections drafted", True),
    ("4 figures inserted (Fig. 1–4)", True),
    ("Corrections applied (§4.4, §4.5, refs)", True),
    ("Paste final abstract", False),
    ("Final read-through + Fig check", False),
    ("Author info from supervisor", False),
]
for i, (item, done) in enumerate(final_items):
    mark = "✅" if done else "⬜"
    col  = BLACK if done else ORANGE
    txbox(s, mark + "  " + item, Inches(0.55), Inches(2.82 + i * 0.65),
          Inches(5.7), Inches(0.6), size=14, color=col)

# Publication column
rect(s, Inches(6.85), Inches(1.4), Inches(6.1), Inches(5.72), BGWHITE, ORANGE, Pt(1.5))
txbox(s, "Publication Paper (KIIT)", Inches(7.05), Inches(1.5), Inches(5.7), Inches(0.5),
      size=18, bold=True, color=ORANGE)
txbox(s, "docs/manuscript-publication.hwp", Inches(7.05), Inches(1.98),
      Inches(5.7), Inches(0.38), size=11, color=GRAY, italic=True)
txbox(s, "Awaiting experimental results (RQ1-error, RQ2)", Inches(7.05), Inches(2.34),
      Inches(5.7), Inches(0.38), size=13, bold=True, color=ORANGE)

pub_items = [
    ("Text corrections (same as final exam)", True),
    ("RQ1 union results", True),
    ("RQ3 ordering results", True),
    ("RQ1 error corpus (108 payloads)", False),
    ("RQ2 Safeline CE transfer eval", False),
    ("Integrate all results + submit", False),
]
for i, (item, done) in enumerate(pub_items):
    mark = "✅" if done else "⬜"
    col  = BLACK if done else ORANGE
    txbox(s, mark + "  " + item, Inches(7.05), Inches(2.82 + i * 0.65),
          Inches(5.7), Inches(0.6), size=14, color=col)


# ─── SLIDE 8: NEXT STEPS (RQ2) ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
hbar(s, "Next Steps — RQ2: Transfer to Safeline CE",
     "Zero-shot policy transfer: ModSecurity (rule-based) → Safeline CE (ML-based WAF)")

rect(s, Inches(0.35), Inches(1.42), Inches(5.6), Inches(5.7),
     RGBColor(0xed, 0xf7, 0xed), GREEN, Pt(1.5))
txbox(s, "Safeline CE Setup", Inches(0.55), Inches(1.52), Inches(5.2), Inches(0.5),
      size=16, bold=True, color=GREEN)
setup_items = [
    "✅  Docker container running",
    "✅  Port 8888 active",
    "✅  Verified: normal=200, SQLi=403",
    "✅  Eval commands prepared",
    "⬜  Run FNR₀ baseline",
    "⬜  Run TRPO/PPO/A2C zero-shot eval",
    "⬜  Compute transferability gap",
]
for i, t in enumerate(setup_items):
    col = BLACK if "✅" in t else ORANGE
    txbox(s, t, Inches(0.55), Inches(2.12 + i * 0.72), Inches(5.2), Inches(0.65),
          size=14, color=col, bold=("⬜" in t))

rect(s, Inches(6.2), Inches(1.42), Inches(6.75), Inches(5.7),
     BGWHITE, BLUE_D, Pt(1.5))
txbox(s, "Research Design", Inches(6.4), Inches(1.52), Inches(6.35), Inches(0.5),
      size=16, bold=True, color=BLUE_D)
design_items = [
    ("Target WAF", "Safeline CE (Chaitin Tech)\nNeural-network based detection"),
    ("Method", "Zero-shot transfer — same trained\nmodel, new WAF endpoint"),
    ("Corpus", "108 union payloads (same as RQ1)"),
    ("Metrics", "IFNR + SPBARC at Safeline"),
    ("Key metric", "Transferability gap\n= IFNR(ModSec) − IFNR(Safeline)"),
]
for i, (lbl, val) in enumerate(design_items):
    txbox(s, lbl + ":", Inches(6.4), Inches(2.12 + i * 1.0), Inches(1.9), Inches(0.85),
          size=13, bold=True, color=BLUE_D)
    txbox(s, val, Inches(8.3), Inches(2.12 + i * 1.0), Inches(4.45), Inches(0.85),
          size=13, color=BLACK)

txbox(s, "Hypothesis: mutations effective against rule-based WAF may be less effective against\n"
         "ML-based WAF — different detection mechanism (token distribution vs. pattern matching)",
      Inches(0.35), Inches(7.15), Inches(12.6), Inches(0.28),
      size=11, color=BLUE_D, italic=True)


# ─── SLIDE 9: TIMELINE ───────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
hbar(s, "Timeline & Milestones")

milestones = [
    ("Done", "Training & eval\nRQ1 + RQ3",         GREEN,  True),
    ("Done", "Figure generation\n+ paper draft",    GREEN,  True),
    ("Now",  "Final exam\npaper finalize",          ORANGE, False),
    ("~Jun 16", "Final exam\nsubmission",           RED,    False),
    ("TBD",  "Run RQ2 eval\n(Safeline)",            BLUE_D, False),
    ("TBD",  "Publication\npaper complete",         BLUE_D, False),
]

bar_y   = Inches(2.0)
bar_h   = Inches(0.08)
total_w = Inches(12.0)
x_start = Inches(0.65)
step_w  = total_w / (len(milestones) - 1)

rect(s, x_start, bar_y, total_w, bar_h, LGRAY)

for i, (date, label, col, done) in enumerate(milestones):
    cx = x_start + Inches(i * (12.0 / (len(milestones) - 1)))
    dot_r = Inches(0.22)
    dot_color = col if done else BGWHITE
    dot = s.shapes.add_shape(9, cx - dot_r, bar_y - dot_r + bar_h / 2, dot_r * 2, dot_r * 2)
    dot.fill.solid()
    dot.fill.fore_color.rgb = dot_color
    dot.line.color.rgb = col
    dot.line.width = Pt(2.5)

    txbox(s, date, cx - Inches(0.85), bar_y - Inches(0.65), Inches(1.7), Inches(0.45),
          size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
    txbox(s, label, cx - Inches(0.9), bar_y + Inches(0.35), Inches(1.8), Inches(0.85),
          size=12, color=col if not done else BLACK, align=PP_ALIGN.CENTER)

# Progress bar fill (first 2 milestones = done)
done_w = x_start + Inches(1 * (12.0 / (len(milestones) - 1)))
rect(s, x_start, bar_y, done_w - x_start, bar_h, GREEN)

rect(s, Inches(0.35), Inches(3.5), Inches(12.6), Inches(3.5), BGWHITE, BLUE_D, Pt(1))
txbox(s, "Immediate priorities:", Inches(0.55), Inches(3.62),
      Inches(12.1), Inches(0.45), size=16, bold=True, color=BLUE_D)
priorities = [
    ("1", "Paste final abstract into final-exam HWP  (text ready in PROJECT_STATE.md)", BLUE_D),
    ("2", "Final read-through — check Fig. 1–4 numbering consistency + typos", BLUE_D),
    ("3", "Run RQ2 Safeline eval (commands ready) → paste output to analyze", ORANGE),
    ("4", "Submit final exam paper by June 16, 2026", RED),
]
for i, (n, task, col) in enumerate(priorities):
    txbox(s, n + ".  " + task, Inches(0.55), Inches(4.15 + i * 0.65),
          Inches(12.1), Inches(0.6), size=14, color=col)


# ─── SLIDE 10: CLOSING ───────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, BLUE_D)
rect(s, 0, Inches(2.5), SW, Inches(2.8), RGBColor(0x19, 0x55, 0x9a))
txbox(s, "Thank You", Inches(1), Inches(1.35), Inches(11.3), Inches(1.0),
      size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, "Discussion & Questions", Inches(1), Inches(2.75), Inches(11.3), Inches(0.8),
      size=28, color=BLUE_L, align=PP_ALIGN.CENTER)

summary_lines = [
    "RQ1 ✅  TRPO 99.1% IFNR  |  PPO 88.9%  |  A2C 76.9%  —  complex tier is the discriminator",
    "RQ3 ✅  35 consensus ordering pairs  —  structural signal from WAF rule topology",
    "Paper ⏳  Final exam near-complete; RQ2 eval pending for publication",
]
for i, line in enumerate(summary_lines):
    txbox(s, line, Inches(1), Inches(4.05 + i * 0.65), Inches(11.3), Inches(0.58),
          size=14, color=RGBColor(0xaa, 0xcc, 0xee), align=PP_ALIGN.CENTER)

txbox(s, "Roby Firnando Yusuf  |  Advisor: Prof. Duksan Ryu",
      Inches(1), Inches(6.65), Inches(11.3), Inches(0.55),
      size=14, color=RGBColor(0x88, 0xaa, 0xcc), align=PP_ALIGN.CENTER)


prs.save(OUT)
print(f"Saved : {OUT}")
print(f"Slides: {len(prs.slides)}")
