"""
Generate SeqSQLi presentation PPT (English).
Run:
    pip install python-pptx
    python3 presentasi/build_ppt.py
Output: presentasi/SeqSQLi_Presentation.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT   = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIGDIR = os.path.join(ROOT, "figures")
OUT    = os.path.join(ROOT, "presentasi", "SeqSQLi_Presentation.pptx")

BLUE_D  = RGBColor(0x21,0x66,0xac)
BLUE_L  = RGBColor(0xdc,0xe8,0xf5)
ORANGE  = RGBColor(0xe0,0x8d,0x00)
RED     = RGBColor(0xd6,0x60,0x4d)
GREEN   = RGBColor(0x4d,0xac,0x26)
WHITE   = RGBColor(0xFF,0xFF,0xFF)
BLACK   = RGBColor(0x1a,0x1a,0x1a)
GRAY    = RGBColor(0x55,0x55,0x55)
BGWHITE = RGBColor(0xF7,0xF9,0xFC)

SW = Inches(13.33)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

def bg(slide, color=BGWHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def txbox(slide, text, x, y, w, h, size=18, bold=False,
          color=BLACK, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color

def rect(slide, x, y, w, h, fill, ec=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if ec:
        s.line.color.rgb = ec
        s.line.width = lw
    else:
        s.line.fill.background()
    return s

def hbar(slide, title, subtitle=None):
    rect(slide, 0, 0, SW, Inches(1.3), BLUE_D)
    txbox(slide, title, Inches(0.45), Inches(0.1),
          Inches(12.4), Inches(0.75), size=30, bold=True, color=WHITE)
    if subtitle:
        txbox(slide, subtitle, Inches(0.45), Inches(0.8),
              Inches(12.4), Inches(0.42), size=15, color=BLUE_L)

def pic(slide, path, x, y, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, w, h)
    else:
        txbox(slide, "[missing: " + os.path.basename(path) + "]",
              x, y, w, h, size=12, color=RED)

def note(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# SLIDE 1 — TITLE
s = prs.slides.add_slide(BLANK)
bg(s, BLUE_D)
rect(s, 0, Inches(2.4), SW, Inches(3.3), RGBColor(0x19,0x55,0x9a))
txbox(s, "SeqSQLi", Inches(1), Inches(1.0), Inches(11.3), Inches(1.1),
      size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, "Sequential SQL Injection WAF Bypass\nvia Deep Reinforcement Learning",
      Inches(1), Inches(2.1), Inches(11.3), Inches(1.2),
      size=24, color=BLUE_L, align=PP_ALIGN.CENTER)
txbox(s, "Roby Firnando Yusuf  |  Advisor: Prof. Duksan Ryu",
      Inches(1), Inches(3.55), Inches(11.3), Inches(0.55),
      size=17, color=RGBColor(0xaa,0xcc,0xee), align=PP_ALIGN.CENTER)
txbox(s, "Final Seminar  |  June 2026",
      Inches(1), Inches(4.1), Inches(11.3), Inches(0.5),
      size=15, color=RGBColor(0x88,0xaa,0xcc), align=PP_ALIGN.CENTER)
note(s, "OPENING (~30 sec)\nGood morning. I am Roby Firnando Yusuf.\nToday I present SeqSQLi, a Deep RL framework for automated WAF bypass.\nLet us get started.")

# SLIDE 2 — AGENDA
s = prs.slides.add_slide(BLANK)
bg(s)
hbar(s, "Agenda")
for i, txt in enumerate([
    "1.  Background & Motivation",
    "2.  SeqSQLi Framework",
    "3.  Experimental Setup",
    "4.  Results & Analysis",
    "5.  Conclusion & Future Work",
]):
    txbox(s, txt, Inches(2.5), Inches(1.65+i*0.9), Inches(9), Inches(0.75),
          size=24, bold=True, color=BLUE_D)
note(s, "AGENDA (~15 sec)\nFive parts: background, framework, setup, results, conclusion.")

# SLIDE 3 — BACKGROUND
s = prs.slides.add_slide(BLANK)
bg(s)
hbar(s, "Background", "SQL Injection remains a persistent web security threat")
items = [
    ("SQL Injection listed in OWASP Top 10 2025 — Injection category (A05)", BLACK),
    ("Attackers embed SQL commands into user-controlled inputs (e.g. login forms)", BLACK),
    ("Impact: unauthorized access, full database exfiltration", BLACK),
    ("Web Application Firewalls (WAFs) are the primary defense layer", BLUE_D),
    ("WAFs block requests by matching text against known malicious patterns", BLACK),
    ("Core weakness: WAFs check syntax, not semantics", RED),
    ("  Same SQL meaning + different surface form = WAF evasion", RED),
]
for i, (txt, col) in enumerate(items):
    prefix = "  " if txt.startswith(" ") else "  -  "
    txbox(s, prefix+txt, Inches(0.5), Inches(1.45+i*0.56),
          Inches(12.3), Inches(0.52), size=17, color=col, bold=(col==BLUE_D))
note(s, "BACKGROUND (~1 min)\nSQL Injection: attacker inserts SQL into inputs. WAFs scan text patterns.\nWeakness: change the form, keep the meaning. Example: SELECT -> SeLeCt.\nSame database execution, WAF pattern misses.")

# SLIDE 4 — MOTIVATION
s = prs.slides.add_slide(BLANK)
bg(s)
hbar(s, "Problem & Motivation", "Why reinforcement learning?")
rect(s, Inches(0.35), Inches(1.42), Inches(5.95), Inches(5.6),
     RGBColor(0xff,0xf3,0xd6), ORANGE, Pt(1.5))
txbox(s, "Prior Approaches",
      Inches(0.55), Inches(1.52), Inches(5.5), Inches(0.5),
      size=17, bold=True, color=ORANGE)
for i, t in enumerate([
    "Random mutation — inefficient",
    "Manual trial-and-error — slow",
    "Grammar-based — no live WAF feedback",
    "No study of mutation ordering",
]):
    txbox(s, "-  "+t, Inches(0.6), Inches(2.15+i*0.82),
          Inches(5.5), Inches(0.72), size=16, color=BLACK)
rect(s, Inches(6.85), Inches(1.42), Inches(5.95), Inches(5.6),
     BLUE_L, BLUE_D, Pt(1.5))
txbox(s, "SeqSQLi (Our Approach)",
      Inches(7.05), Inches(1.52), Inches(5.5), Inches(0.5),
      size=17, bold=True, color=BLUE_D)
for i, t in enumerate([
    "RL agent learns from live WAF responses",
    "Discovers optimal mutation sequences",
    "Trained online against a real WAF",
    "Formally measures mutation ordering effect",
]):
    txbox(s, "-  "+t, Inches(7.1), Inches(2.15+i*0.82),
          Inches(5.5), Inches(0.72), size=16, color=BLACK)
note(s, "MOTIVATION (~1 min)\nPrior work: random/grammar-based, no ordering study.\nSeqSQLi: RL agent gets direct WAF feedback, learns optimal sequences.\nFirst work to formally study whether mutation order is a learnable signal.")

# SLIDE 5 — FRAMEWORK
s = prs.slides.add_slide(BLANK)
bg(s)
hbar(s, "SeqSQLi Framework", "Modeled as a Markov Decision Process (MDP)")
pic(s, os.path.join(FIGDIR,"fig_architecture.png"),
    Inches(3.9), Inches(1.38), Inches(5.6), Inches(5.7))
txbox(s, "MDP Components:", Inches(0.3), Inches(1.5), Inches(3.4), Inches(0.45),
      size=15, bold=True, color=BLUE_D)
for i, (label, val, col) in enumerate([
    ("State  (s)",   "67-dim observation vector",  BLUE_D),
    ("Action  (a)",  "51 mutation operators",       ORANGE),
    ("Reward  (r)",  "WAF response -> score",       RED),
    ("Environment",  "Live ModSecurity WAF",        GREEN),
    ("Episode",      "Max 15 mutation steps",       GRAY),
]):
    txbox(s, label, Inches(0.35), Inches(2.05+i*0.78), Inches(1.6), Inches(0.6),
          size=13, bold=True, color=col)
    txbox(s, val, Inches(1.95), Inches(2.05+i*0.78), Inches(1.9), Inches(0.6),
          size=13, color=BLACK)
note(s, "FRAMEWORK (~1.5 min)\nMDP: observe 67-dim state -> select mutation -> HTTP request to WAF -> response -> reward -> update policy.\nEpisode: one payload, max 15 steps, ends on SUCCESS or truncation.")

# SLIDE 6 — STATE
s = prs.slides.add_slide(BLANK)
bg(s)
hbar(s, "State Representation", "Agent observes the payload as a 67-dimensional vector")
boxes = [
    ("1. Payload Features (14 dim)",
     "Binary flags: which mutation families\nhave already been applied",
     BLUE_D, BLUE_L),
    ("2. Injection Type Bit (1 dim)",
     "0 = union-based,  1 = error-based\nPrevents policy collapse across tasks",
     BLUE_D, BLUE_L),
    ("3. Last-Action One-Hot (51 dim)",
     "Which mutation was applied at t-1\nEnables learning of ordering dependencies",
     ORANGE, RGBColor(0xff,0xf3,0xd6)),
    ("4. Step Counter (1 dim)",
     "Current step normalized by Tmax\nImplicit remaining-budget signal",
     GREEN, RGBColor(0xed,0xf7,0xed)),
]
for i, (title, body, ec, fc) in enumerate(boxes):
    c, r = i%2, i//2
    bx = Inches(0.35+c*6.45)
    by = Inches(1.45+r*2.75)
    rect(s, bx, by, Inches(6.1), Inches(2.55), fc, ec, Pt(1.5))
    txbox(s, title, bx+Inches(0.15), by+Inches(0.1), Inches(5.8), Inches(0.8),
          size=16, bold=True, color=ec)
    txbox(s, body, bx+Inches(0.15), by+Inches(0.85), Inches(5.8), Inches(1.5),
          size=14, color=BLACK)
txbox(s, "Total: 14 + 1 + 51 + 1 = 67 dimensions",
      Inches(0.35), Inches(7.05), Inches(12.6), Inches(0.35),
      size=15, bold=True, color=BLUE_D, align=PP_ALIGN.CENTER)
note(s, "STATE (~1 min)\n14: checklist of applied mutation families.\n1: task type bit prevents policy collapse.\n51: last-action one-hot — key for ordering learning.\n1: step fraction — encourages efficiency.")

# SLIDE 7 — ACTION SPACE
s = prs.slides.add_slide(BLANK)
bg(s)
hbar(s, "Action Space",
     "51 grammar-aware mutation operators across 9 families — semantics preserved")
families = [
    ("Case Mixing",          "SeLeCt,  UnIoN",          BLUE_D),
    ("Whitespace Variants",  "tab,  newline,  %0a",      BLUE_D),
    ("Comment Insertion",    "/**/,  /*!50000*/",        BLUE_D),
    ("Hex / Encoding",       "0x3a,  URL-encode",        ORANGE),
    ("Identifier Tricks",    "`users`,  double-write",   ORANGE),
    ("Aggregate Swap",       "JSON_ARRAYAGG(CONCAT(.))", ORANGE),
    ("Null Byte",            "%00,  \\x00",              RED),
    ("Paren / Space",        "func( ),  dot prefix",     RED),
    ("Semantic Substitution","AND->&&,  OR->||",         GREEN),
]
for i, (fam, ex, col) in enumerate(families):
    r, c = divmod(i, 3)
    bx = Inches(0.35+c*4.3)
    by = Inches(1.42+r*1.8)
    rect(s, bx, by, Inches(4.05), Inches(1.6), BGWHITE, col, Pt(1.2))
    txbox(s, fam, bx+Inches(0.1), by+Inches(0.1), Inches(3.85), Inches(0.6),
          size=14, bold=True, color=col)
    txbox(s, ex,  bx+Inches(0.1), by+Inches(0.7), Inches(3.85), Inches(0.75),
          size=13, color=GRAY, italic=True)
note(s, "ACTIONS (~1 min)\n51 operators, 9 families. All preserve SQL semantics.\nKey: aggregate swap (GROUP_CONCAT -> JSON_ARRAYAGG) unlocked complex tier bypass.")

# SLIDE 8 — REWARD
s = prs.slides.add_slide(BLANK)
bg(s)
hbar(s, "Reward Function", "Feedback signal guiding the agent's learning")
for i, (outcome, score, desc, col) in enumerate([
    ("SUCCESS",      "+10.0", "WAF bypassed — canary markers appear in response body",  GREEN),
    ("SQL_ERROR",    "+0.5",  "Payload reached DB but caused a SQL syntax error",        ORANGE),
    ("FILTERED",     "-1.0",  "WAF blocked with a specific rule match",                  RED),
    ("UNKNOWN",      "-0.5",  "Unrecognized response pattern",                           GRAY),
    ("SERVER_ERROR", "-1.5",  "HTTP 500 — server error",                                RED),
    ("WAF_BLOCKED",  "-2.0",  "WAF hard-blocked the request (HTTP 403)",                RED),
    ("STAGNANT",     "-3.0",  "Mutation produced no change to the payload",              RED),
]):
    rect(s, Inches(0.35), Inches(1.38+i*0.73), Inches(12.65), Inches(0.69),
         RGBColor(0xf7,0xf9,0xfc) if i%2==0 else WHITE)
    txbox(s, outcome, Inches(0.5),  Inches(1.42+i*0.73), Inches(2.2), Inches(0.55),
          size=14, bold=True, color=col)
    txbox(s, score,   Inches(2.75), Inches(1.42+i*0.73), Inches(1.0), Inches(0.55),
          size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    txbox(s, desc,    Inches(3.9),  Inches(1.42+i*0.73), Inches(9.0), Inches(0.55),
          size=13, color=BLACK)
txbox(s, "+ PBRS Shaping: bonus per WAF trigger removed — policy-invariant dense guidance (Ng et al. 1999)",
      Inches(0.35), Inches(6.7), Inches(12.65), Inches(0.4),
      size=12, color=BLUE_D, italic=True)
note(s, "REWARD (~45 sec)\nSeven outcomes +10 to -3. STAGNANT lowest to prevent no-op policy.\nPBRS: bonus per WAF trigger removed, dense guidance, does not change optimal policy.")

# SLIDE 9 — ALGORITHMS
s = prs.slides.add_slide(BLANK)
bg(s)
hbar(s, "Three RL Algorithms Compared",
     "All on-policy, trained under identical conditions")
for i, (abbr, full, desc, fc, ec) in enumerate([
    ("PPO",  "Proximal Policy Optimization",
     "Clips the policy update ratio\nto prevent excessively large steps.\nSoft brake on policy change.",
     RGBColor(0xff,0xf3,0xd6), ORANGE),
    ("TRPO", "Trust Region Policy Optimization",
     "Enforces a hard KL-divergence\nconstraint on each update.\nGuarantees monotonic improvement.",
     BLUE_L, BLUE_D),
    ("A2C",  "Advantage Actor-Critic",
     "No update constraint — gradients\napplied in a single direct pass.\nSimplest but most noise-sensitive.",
     RGBColor(0xfd,0xe0,0xdc), RED),
]):
    bx = Inches(0.35+i*4.3)
    rect(s, bx, Inches(1.4), Inches(4.15), Inches(5.65), fc, ec, Pt(1.5))
    txbox(s, abbr, bx+Inches(0.15), Inches(1.5), Inches(3.85), Inches(0.8),
          size=30, bold=True, color=ec, align=PP_ALIGN.CENTER)
    txbox(s, full, bx+Inches(0.15), Inches(2.25), Inches(3.85), Inches(0.6),
          size=12, color=GRAY, align=PP_ALIGN.CENTER, italic=True)
    txbox(s, desc, bx+Inches(0.15), Inches(2.95), Inches(3.85), Inches(2.8),
          size=15, color=BLACK)
txbox(s, "Shared: two-layer MLP (64 units, ReLU)  |  150,000 timesteps  |  GAE (lam=0.95)  |  gamma=0.99",
      Inches(0.35), Inches(7.1), Inches(12.65), Inches(0.35),
      size=13, color=GRAY, align=PP_ALIGN.CENTER, italic=True)
note(s, "ALGORITHMS (~1.5 min)\nPPO: clipped ratio, soft constraint.\nTRPO: hard KL bound, monotonic improvement guarantee — suits sparse/long-horizon.\nA2C: no constraint, fast but unstable.\nSame network, budget, environment — difference = update mechanism only.")

# SLIDE 10 — SETUP
s = prs.slides.add_slide(BLANK)
bg(s)
hbar(s, "Experimental Setup", "Identical conditions for all three algorithms")
for i, (lbl, val) in enumerate([
    ("WAF Target",   "ModSecurity CRS v3.3.2\nnginx 1.18  |  sqli-labs Less-1"),
    ("URL",          "http://localhost:8080/Less-1/"),
    ("Hardware",     "Ubuntu 24.04 (ARM64)\n121 GiB RAM  |  NVIDIA GB10 GPU"),
    ("Inference",    "CPU-only (lightweight MLP)"),
]):
    rect(s, Inches(0.35), Inches(1.42+i*1.42), Inches(6.1), Inches(1.28),
         BGWHITE, BLUE_D, Pt(1))
    txbox(s, lbl, Inches(0.5),  Inches(1.49+i*1.42), Inches(1.75), Inches(0.5),
          size=13, bold=True, color=BLUE_D)
    txbox(s, val, Inches(2.3),  Inches(1.49+i*1.42), Inches(4.0),  Inches(1.05),
          size=13, color=BLACK)
for i, (lbl, val) in enumerate([
    ("Payload Corpus","108 validated union-based payloads\n3 complexity tiers (36 each)"),
    ("Tiers",         "Trivial: case + whitespace\nMedium: + aggregate rule\nComplex: + hex + FROM users"),
    ("Training",      "150,000 timesteps  |  max 15 steps/ep"),
    ("Evaluation",    "Stochastic mode  |  all 108 payloads\nMetrics: IFNR and SPBARC"),
]):
    rect(s, Inches(6.9),  Inches(1.42+i*1.42), Inches(6.1), Inches(1.28),
         BGWHITE, ORANGE, Pt(1))
    txbox(s, lbl, Inches(7.05), Inches(1.49+i*1.42), Inches(1.75), Inches(0.5),
          size=13, bold=True, color=ORANGE)
    txbox(s, val, Inches(8.85), Inches(1.49+i*1.42), Inches(4.0),  Inches(1.05),
          size=13, color=BLACK)
note(s, "SETUP (~45 sec)\nAll algorithms: same environment, corpus, budget, eval protocol.\nTier = number of WAF rule families to defeat simultaneously.\nAll 108 payloads semantically validated against raw backend before training.")

# SLIDE 11 — RESULTS TABLE
s = prs.slides.add_slide(BLANK)
bg(s)
hbar(s, "RQ1 — Overall WAF Bypass Performance",
     "TRPO achieves the best results across all metrics")
col_hdrs = ["Algorithm","IFNR (higher)","SPBARC (lower)","Trivial","Medium","Complex"]
col_ws   = [Inches(2.2), Inches(1.85), Inches(1.95), Inches(1.55), Inches(1.55), Inches(1.6)]
x0, y0  = Inches(0.5), Inches(1.5)
x = x0
for hdr, w in zip(col_hdrs, col_ws):
    rect(s, x, y0, w, Inches(0.58), BLUE_D)
    txbox(s, hdr, x+Inches(0.05), y0+Inches(0.06), w-Inches(0.1), Inches(0.46),
          size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    x += w
for i, (algo, ifnr, spbarc, triv, med, comp, ec, fc) in enumerate([
    ("TRPO","99.1%","6.07", "100%","100%","97.2%",BLUE_D, BLUE_L),
    ("PPO", "88.9%","6.89", "100%","100%","66.7%",ORANGE, RGBColor(0xff,0xf9,0xee)),
    ("A2C", "76.9%","10.08","100%","97.2%","33.3%",RED,   RGBColor(0xfd,0xf0,0xee)),
]):
    x = x0
    for j, (val, w) in enumerate(zip([algo,ifnr,spbarc,triv,med,comp], col_ws)):
        rect(s, x, y0+Inches(0.58+i*0.8), w, Inches(0.76), fc)
        txbox(s, val, x+Inches(0.05), y0+Inches(0.66+i*0.8), w-Inches(0.1), Inches(0.6),
              size=17, bold=(j in [0,5]),
              color=ec if j in [0,5] else BLACK, align=PP_ALIGN.CENTER)
        x += w
txbox(s, "IFNR = % of blocked payloads ultimately bypassed  |  SPBARC = avg. HTTP requests per bypass",
      Inches(0.5), Inches(4.7), Inches(12.35), Inches(0.4),
      size=11, color=GRAY, italic=True)
txbox(s, "All algorithms substantially outperform the random baseline  (IFNR 3.7%,  SPBARC 216)",
      Inches(0.5), Inches(5.2), Inches(12.35), Inches(0.45),
      size=15, bold=True, color=BLUE_D)
note(s, "RESULTS TABLE (~1.5 min)\nTRPO: 99.1% IFNR, SPBARC 6.07 — 107/108 payloads bypassed.\nPPO 88.9%, A2C 76.9%. All vastly outperform random (3.7%).\nKey story: trivial/medium saturated, complex tier is the divergence point.")

# SLIDE 12 — TIER CHART
s = prs.slides.add_slide(BLANK)
bg(s)
hbar(s, "Tier-Level Analysis",
     "Complex tier is the key discriminator — longer mutation chains required")
pic(s, os.path.join(FIGDIR,"fig_tier_barchart.png"),
    Inches(0.4), Inches(1.38), Inches(7.9), Inches(5.65))
txbox(s, "Why does Complex tier matter?",
      Inches(8.55), Inches(1.55), Inches(4.45), Inches(0.5),
      size=15, bold=True, color=BLUE_D)
for i, (arrow, txt, col) in enumerate([
    ("-", "Trivial: case + whitespace only",            BLACK),
    ("-", "Medium: +1 rule (aggregate functions)",      BLACK),
    ("-", "Complex: +2 more rules\n(hex + FROM users)", BLACK),
    ("->","Requires 3+ ordered mutations\nto defeat all rules simultaneously", BLUE_D),
    ("->","TRPO KL bound provides stability\nfor long mutation chains",        BLUE_D),
]):
    txbox(s, arrow+"  "+txt, Inches(8.55), Inches(2.15+i*1.05),
          Inches(4.5), Inches(0.95), size=14, color=col, bold=(col==BLUE_D))
note(s, "TIER (~1 min)\nTrivial/medium near 100% for all — easy.\nComplex: TRPO 97.2%, PPO 66.7%, A2C 33.3% — dramatic divergence.\nComplex requires 3+ mutations in correct order against different WAF rule families.")

# SLIDE 13 — TRAINING CURVES
s = prs.slides.add_slide(BLANK)
bg(s)
hbar(s, "Training Convergence",
     "TRPO most stable; A2C suffers policy collapse at ~80k timesteps")
pic(s, os.path.join(FIGDIR,"fig_training_curves.png"),
    Inches(0.8), Inches(1.38), Inches(9.6), Inches(5.65))
for i, (algo, desc, col) in enumerate([
    ("TRPO","Converges to highest reward;\nmost stable trajectory", BLUE_D),
    ("PPO", "Faster initial rise;\nplateaus below TRPO",            ORANGE),
    ("A2C", "Policy collapse at ~80k\n(unconstrained update)",      RED),
]):
    txbox(s, algo+": "+desc, Inches(10.6), Inches(2.7+i*1.4), Inches(2.5), Inches(1.2),
          size=13, color=col, bold=(algo=="A2C"))
note(s, "CURVES (~45 sec)\nTRPO: converges highest, most stable.\nPPO: rises fast, plateaus below TRPO.\nA2C: policy collapse at 80k — unconstrained update destroys learned policy.")

# SLIDE 14 — ORDERING
s = prs.slides.add_slide(BLANK)
bg(s)
hbar(s, "RQ3 — Does Mutation Order Matter?",
     "Yes — 35 consensus pairs succeed only in the correct order")
pic(s, os.path.join(FIGDIR,"fig_ordering_pairs.png"),
    Inches(0.3), Inches(1.38), Inches(8.4), Inches(5.8))
txbox(s, "Key Findings:", Inches(8.95), Inches(1.55), Inches(4.05), Inches(0.45),
      size=15, bold=True, color=BLUE_D)
for i, (txt, col) in enumerate([
    ("35 consensus pairs: A->B succeeds;\nB->A fails",       BLUE_D),
    ("Strongest gap:\nident_backtick -> hex_to_char\nFwd 98%  |  Rev 0%", RED),
    ("Pattern consistent across\nall 3 algorithms",          GREEN),
    ("Ordering originates from WAF\nrule topology, not\nalgorithm dynamics", BLUE_D),
]):
    rect(s, Inches(8.95), Inches(2.1+i*1.25), Inches(4.1), Inches(1.15),
         BGWHITE, col, Pt(1))
    txbox(s, txt, Inches(9.1), Inches(2.17+i*1.25), Inches(3.8), Inches(1.0),
          size=13, color=BLACK)
note(s, "ORDERING (~1.5 min)\n35 consensus pairs from training logs.\nMost extreme: ident_backtick->hex_to_char: 98% vs 0%.\nhex_to_char corrupts identifier token before backtick can be applied.\nSame pattern in 3 algorithms -> originates from WAF rule structure, not learning dynamics.")

# SLIDE 15 — CONCLUSION
s = prs.slides.add_slide(BLANK)
bg(s)
hbar(s, "Conclusion")
for i, (col, title, body) in enumerate([
    (BLUE_D, "1.  SeqSQLi effectively learns to bypass WAF",
     "TRPO achieves 99.1% bypass rate — 107 out of 108 initially blocked payloads succeeded"),
    (BLUE_D, "2.  TRPO > PPO > A2C across all metrics",
     "Trust-region / KL-bound constraint is superior for sparse-reward, long-horizon episodes"),
    (GREEN,  "3.  Mutation ordering is a learnable structural signal",
     "35 cross-algorithm consensus pairs confirm ordering originates from WAF rule topology"),
    (ORANGE, "4.  Complex tier is the true discriminator",
     "Trivial and medium tiers are saturated; complex tier reveals real capability differences"),
]):
    rect(s, Inches(0.4), Inches(1.42+i*1.42), Inches(12.55), Inches(1.28),
         BGWHITE, col, Pt(1.3))
    txbox(s, title, Inches(0.6), Inches(1.49+i*1.42), Inches(12.1), Inches(0.52),
          size=16, bold=True, color=col)
    txbox(s, body,  Inches(0.6), Inches(1.97+i*1.42), Inches(12.1), Inches(0.58),
          size=14, color=BLACK)
note(s, "CONCLUSION (~1 min)\n1. RL works for WAF bypass. TRPO nearly solves the problem.\n2. TRPO wins — KL bound advantage in sparse/long-horizon.\n3. Ordering is real and structural, not an artifact.\n4. Complex tier is the true benchmark.")

# SLIDE 16 — FUTURE WORK
s = prs.slides.add_slide(BLANK)
bg(s)
hbar(s, "Limitations & Future Work")
rect(s, Inches(0.35), Inches(1.4), Inches(6.1), Inches(5.7), BGWHITE, RED, Pt(1.5))
txbox(s, "Limitations", Inches(0.55), Inches(1.5), Inches(5.7), Inches(0.55),
      size=20, bold=True, color=RED)
for i, t in enumerate([
    "Single WAF target (ModSecurity only)",
    "Single injection type (union-based)",
    "Single application & injection point",
    "No evaluation on commercial WAFs",
]):
    txbox(s, "-  "+t, Inches(0.6), Inches(2.2+i*1.1), Inches(5.6), Inches(0.95),
          size=16, color=BLACK)
rect(s, Inches(6.85), Inches(1.4), Inches(6.1), Inches(5.7), BGWHITE, GREEN, Pt(1.5))
txbox(s, "Future Work", Inches(7.05), Inches(1.5), Inches(5.7), Inches(0.55),
      size=20, bold=True, color=GREEN)
for i, t in enumerate([
    "RQ2: Zero-shot transfer to Safeline CE\n(ML-based WAF, neural network detection)",
    "Extend corpus to error-based (216 total)",
    "Explore XSS & second-order SQLi",
    "Publication to KIIT Journal",
]):
    txbox(s, "-  "+t, Inches(7.1), Inches(2.2+i*1.1), Inches(5.6), Inches(0.95),
          size=16, color=BLACK)
note(s, "FUTURE WORK (~45 sec)\nLimitations: one WAF, one injection type, one app.\nMost exciting: RQ2 zero-shot transfer to Safeline CE (ML-based WAF).")

# SLIDE 17 — CLOSING
s = prs.slides.add_slide(BLANK)
bg(s, BLUE_D)
rect(s, 0, Inches(2.7), SW, Inches(2.6), RGBColor(0x19,0x55,0x9a))
txbox(s, "Thank You", Inches(1), Inches(1.45), Inches(11.3), Inches(1.1),
      size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, "Questions?", Inches(1), Inches(2.9), Inches(11.3), Inches(0.9),
      size=32, color=BLUE_L, align=PP_ALIGN.CENTER)
txbox(s, "Roby Firnando Yusuf  |  roby@example.com",
      Inches(1), Inches(4.3), Inches(11.3), Inches(0.55),
      size=16, color=RGBColor(0x88,0xaa,0xcc), align=PP_ALIGN.CENTER)
note(s, "CLOSING\nThank you for your attention. Happy to take questions.")

prs.save(OUT)
print(f"Saved : {OUT}")
print(f"Slides: {len(prs.slides)}")
