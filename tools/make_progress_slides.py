"""
tools/make_progress_slides.py
=============================
Generate a CLEAN/minimal lab-meeting progress deck (.pptx) for SeqSQLi.
Light background, no solid header bars — just a thin accent underline.
Bilingual (ID + EN) speaker notes. No network.
Run: python3 -m tools.make_progress_slides
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK    = RGBColor(0x23, 0x27, 0x2F)   # near-black text
MUT    = RGBColor(0x6B, 0x72, 0x80)   # muted grey
ACCENT = RGBColor(0x3A, 0x6E, 0xA5)   # soft blue
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
RED    = RGBColor(0xB0, 0x3A, 0x2E)
CARD   = RGBColor(0xF4, 0xF6, 0xF9)   # very light card
HEADT  = RGBColor(0xE9, 0xEE, 0xF4)   # very light table header tint
ZEBRA  = RGBColor(0xF7, 0xF9, 0xFB)   # very light zebra row
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _noshadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def header(slide, title, sub=None):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.33), Inches(12.2), Inches(0.95))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = INK
    if sub:
        p2 = tf.add_paragraph(); p2.text = sub
        p2.font.size = Pt(13); p2.font.color.rgb = MUT
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.28), Inches(2.0), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background(); _noshadow(ln)


def bullets(slide, items, left=Inches(0.6), top=Inches(1.55),
            width=Inches(12.1), height=Inches(5.4), size=18):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        txt, lvl, color, bold = it, 0, INK, False
        if isinstance(it, tuple):
            txt = it[0]; lvl = it[1] if len(it) > 1 else 0
            color = it[2] if len(it) > 2 else INK
            bold = it[3] if len(it) > 3 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("•  " + txt) if lvl == 0 else ("–  " + txt)
        p.level = lvl
        p.font.size = Pt(size - 2 * lvl); p.font.color.rgb = color; p.font.bold = bold
        p.space_after = Pt(7)
    return box


def table(slide, data, left, top, width, height, font=13, col_widths=None,
          highlight_rows=None):
    rows, cols = len(data), len(data[0])
    g = slide.shapes.add_table(rows, cols, left, top, width, height).table
    # neutralize the default banded style by explicit fills
    if col_widths:
        for c, w in enumerate(col_widths):
            g.columns[c].width = w
    highlight_rows = highlight_rows or {}
    for r in range(rows):
        for c in range(cols):
            cell = g.cell(r, c)
            cell.text = str(data[r][c])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
            cell.margin_left = Pt(7)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(font); para.font.color.rgb = INK
            para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = HEADT; para.font.bold = True
            elif r in highlight_rows:
                cell.fill.fore_color.rgb = highlight_rows[r]
            else:
                cell.fill.fore_color.rgb = ZEBRA if r % 2 else WHITE
    return g


def card(slide, left, top, width, height, fill=CARD):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.color.rgb = RGBColor(0xDD, 0xE3, 0xEA); c.line.width = Pt(0.75)
    _noshadow(c)
    return c


def card_text(slide, left, top, width, height, lines):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (txt, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = bold
        p.space_after = Pt(3)
    return tb


def footnote(slide, text, color=GREEN, size=15, top=Inches(6.7)):
    tb = slide.shapes.add_textbox(Inches(0.6), top, Inches(12.1), Inches(0.7))
    p = tb.text_frame.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.bold = True; p.font.color.rgb = color
    tb.text_frame.word_wrap = True


# =============================================================== 1 Title
s = prs.slides.add_slide(BLANK)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.55), Inches(3.0), Pt(4))
ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background(); _noshadow(ln)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.75), Inches(11.5), Inches(3))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "SeqSQLi"
p.font.size = Pt(52); p.font.bold = True; p.font.color.rgb = INK
p2 = tf.add_paragraph(); p2.text = "Reinforcement Learning for Sequential WAF Evasion"
p2.font.size = Pt(24); p2.font.color.rgb = ACCENT
p3 = tf.add_paragraph(); p3.text = "Lab Meeting — Progress Report   ·   2026-06-19"
p3.font.size = Pt(16); p3.font.color.rgb = MUT

# =============================================================== 2 Recap
s = prs.slides.add_slide(BLANK)
header(s, "Recap — Problem & Approach")
bullets(s, [
    ("WAFs block attacks by pattern-matching; a payload's surface can be changed", 0),
    ("without changing its attack meaning.", 0),
    ("SeqSQLi = an RL agent that applies ONE mutation per step (a sequence).", 0, ACCENT, True),
    ("Reward = a real bypass: request passes the WAF (200) AND extracts the data marker.", 1, MUT),
    ("Three on-policy algorithms compared: PPO, TRPO, A2C.", 0),
    ("Two WAFs: ModSecurity (rule-based) and Chaitin Safeline (ML-based), live in Docker.", 0),
], size=19)

# =============================================================== 3 RQs
s = prs.slides.add_slide(BLANK)
header(s, "Research Questions — Status")
table(s, [
    ["RQ", "Question", "Status"],
    ["RQ1", "Which deep-RL algorithm is most effective & efficient?", "DONE (union + error)"],
    ["RQ2", "Does a rule-based-trained policy transfer to an ML-based WAF?", "DONE — key result"],
    ["RQ3", "Does mutation ORDER affect success?", "DONE"],
], Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.3), font=15,
   col_widths=[Inches(1.1), Inches(7.9), Inches(3.1)])
footnote(s, "New this week: RQ1 error corpus (multi-seed) + RQ2 Safeline transfer.",
         color=ACCENT, top=Inches(4.3))

# =============================================================== 4 Metrics (NEW)
s = prs.slides.add_slide(BLANK)
header(s, "Metrics — two numbers, in plain words")
card(s, Inches(0.55), Inches(1.55), Inches(6.0), Inches(4.55))
card_text(s, Inches(0.8), Inches(1.75), Inches(5.5), Inches(4.2), [
    ("IFNR", 24, ACCENT, True),
    ("Induced False Negative Rate", 13, MUT, False),
    ("How OFTEN the agent succeeds.  Higher = better.", 13, INK, True),
    ("", 8, INK, False),
    ("FN  = WAF fails to block (attack passes)", 12, INK, False),
    ("FNR₀ = passes WITHOUT agent = 0% here", 12, INK, False),
    ("MFNR = passes WITH the agent's mutations", 12, INK, False),
    ("IFNR = MFNR − FNR₀  (the increase the agent causes)", 12, INK, True),
    ("", 8, INK, False),
    ("e.g. 99% = broke through 99% of payloads the WAF blocked", 12, GREEN, False),
])
card(s, Inches(6.78), Inches(1.55), Inches(6.0), Inches(4.55))
card_text(s, Inches(7.03), Inches(1.75), Inches(5.5), Inches(4.2), [
    ("SPBARC", 24, ACCENT, True),
    ("Successful Payload Bypass Average Request Count", 12, MUT, False),
    ("How EFFICIENT the agent is.  Lower = better.", 13, INK, True),
    ("", 8, INK, False),
    ("= total requests  ÷  number of successful bypasses", 12, INK, True),
    ("(average tries needed for one bypass)", 12, MUT, False),
    ("", 8, INK, False),
    ("e.g. TRPO 6  → ~6 tries per bypass (efficient)", 12, INK, False),
    ("      A2C 10 → ~10 tries", 12, INK, False),
    ("      Random 216 → very wasteful", 12, RED, False),
])
footnote(s, "Ideal agent = high IFNR + low SPBARC.  TRPO wins both (99.1% / 6.07).",
         color=GREEN, top=Inches(6.35))

# =============================================================== 5 RQ1 union
s = prs.slides.add_slide(BLANK)
header(s, "RQ1 — Algorithm Comparison (Union corpus, ModSecurity)",
       "108 validated payloads · 150k steps · identical environment")
table(s, [
    ["Algorithm", "IFNR", "SPBARC", "trivial", "medium", "complex"],
    ["TRPO", "+99.1%", "6.07", "100%", "100%", "97.2%"],
    ["PPO",  "+88.9%", "6.89", "100%", "100%", "66.7%"],
    ["A2C",  "+76.9%", "10.08", "100%", "97.2%", "33.3%"],
    ["Random (baseline)", "+3.7%", "216", "0%", "—", "—"],
], Inches(0.8), Inches(1.9), Inches(11.7), Inches(2.5), font=15,
   highlight_rows={1: RGBColor(0xE5, 0xF1, 0xE8)})
bullets(s, [
    ("Ranking TRPO > PPO > A2C. Difference shows up in the COMPLEX tier (longest chains).", 0, GREEN, True),
    ("Trust-region / clipped methods beat vanilla A2C under sparse, long-horizon reward.", 0, MUT),
], top=Inches(4.8), size=17)

# =============================================================== 6 RQ1 error
s = prs.slides.add_slide(BLANK)
header(s, "RQ1 — Error corpus (NEW, multi-seed)",
       "108 error-based payloads · same pipeline · 3 random seeds each")
table(s, [
    ["Algorithm", "IFNR (mean ± std)", "range (3 seeds)", "vs union"],
    ["TRPO", "30.3% ± 2.3%", "27.8 – 32.4", "(was 99.1%)"],
    ["PPO",  "26.3% ± 1.9%", "24.1 – 27.8", "(was 88.9%)"],
    ["A2C",  "22.3% ± 14.4%", "5.6 – 30.6", "(was 76.9%)"],
], Inches(0.8), Inches(1.9), Inches(11.7), Inches(2.2), font=15,
   highlight_rows={3: RGBColor(0xF7, 0xE7, 0xE5)})
bullets(s, [
    ("Error payloads are MUCH harder: best IFNR drops 99% → ~30%.", 0, RED, True),
    ("A2C is UNSTABLE: std ≈ 14% (~7× TRPO/PPO); swings 5.6%..30.6% by seed only.", 0, RED, True),
    ("Running 3 seeds mattered: one A2C seed was 0% — reporting it alone = a wrong claim.", 0, MUT),
], top=Inches(4.4), size=17)

# =============================================================== 7 RQ3
s = prs.slides.add_slide(BLANK)
header(s, "RQ3 — Mutation Order Matters")
bullets(s, [
    ("Same mutations, different ORDER → very different success.", 0),
    ("35 cross-algorithm consensus pairs: reversing the order sharply drops success.", 0, ACCENT, True),
    ("Precondition example: ident_backtick → hex_to_char works;  reversed = 0%.", 1, MUT),
    ("Ordering is a STRUCTURAL / causal signal from the WAF rule topology,", 0),
    ("not an artifact of any single algorithm.", 0),
], size=19)

# =============================================================== 8 RQ2 transfer
s = prs.slides.add_slide(BLANK)
header(s, "RQ2 — Transfer: ModSecurity (regex) → Safeline (ML)",
       "Zero-shot: take ModSec-trained models, test on Safeline (no retraining)")
bullets(s, [
    ("Transfer IFNR = 0% — ALL algorithms, BOTH corpora (12 models). Every mutated", 0, RED, True),
    ("payload (~900) hit 403. The agent did not give up — it genuinely cannot pass.", 1, MUT),
    ("Transferability gap ≈ the entire ModSec IFNR (e.g. 99.1% → 0%).", 0, ACCENT, True),
    ("0% is GENUINE, not a bug — positive control passes:", 0),
    ("benign → 200;  textbook SQLi → 403 (WAF active);  crafted PoC → 200 + data.", 1, GREEN),
], size=18)

# =============================================================== 9 three-regime
s = prs.slides.add_slide(BLANK)
header(s, "RQ2 — Training directly on Safeline: the key finding")
table(s, [
    ["Setting", "WAF type", "Passes WAF", "Steals data (IFNR)"],
    ["ModSecurity", "regex", "high", "76 – 99%"],
    ["Safeline — zero-shot transfer", "ML", "~0% (all 403)", "0%"],
    ["Safeline — trained directly", "ML", "~89%", "0%"],
], Inches(0.7), Inches(1.85), Inches(11.9), Inches(2.3), font=15,
   highlight_rows={3: RGBColor(0xFB, 0xF1, 0xDA)},
   col_widths=[Inches(4.1), Inches(2.0), Inches(2.8), Inches(3.0)])
bullets(s, [
    ("On regex-WAF: passing & stealing data go TOGETHER. On ML-WAF: they SPLIT.", 0, ACCENT, True),
    ("Trained directly, the agent passes Safeline ~89% — but steals data 0%.", 0),
    ("846 / 972 requests = HTTP 200 but SQL_ERROR → it passed only by breaking the SQL.", 0, RED, True),
], top=Inches(4.35), size=17)

# =============================================================== 10 mechanism
s = prs.slides.add_slide(BLANK)
header(s, "RQ2 — Why it 'passes the WAF but steals nothing'")
bullets(s, [
    ("The ML detector judges meaning, not surface patterns.", 0),
    ("To look harmless, the agent converges to one 'break-everything' chain:", 0),
    ("double-URL-encode → UNION is no longer a keyword", 1, RED),
    ("GBK trick → breaks the quote on a UTF-8 backend", 1, RED),
    ("Result: neither Safeline NOR MySQL sees valid SQL → passes (200) but errors (no data).", 0, ACCENT, True),
    ("Takeaway: on ML-WAFs the agent can only evade by sacrificing the attack itself.", 0, GREEN, True),
], size=18)

# =============================================================== 11 novelty
s = prs.slides.add_slide(BLANK)
header(s, "Positioning & Novelty (vs BWAFSQLi, ACM TOSEM 2026)")
bullets(s, [
    ("BWAFSQLi already bypasses Safeline (semantic-preserving) and already beats RL methods.", 0),
    ("So 'RL finds a Safeline bypass' is NOT our contribution.", 1, MUT),
    ("What is genuinely new (prior work does not measure this):", 0, ACCENT, True),
    ("(1) Transferability gap rule-based → ML-based WAF, in IFNR / SPBARC.", 1),
    ("(2) Causal mutation-ORDER analysis (structural, not algorithm-specific).", 1),
    ("(3) 'Evade vs steal-data' decoupling on ML-WAFs, with a concrete mechanism.", 1),
    ("Honest framing: we study WHY structural policies transfer — not a new bypass tool.", 0, GREEN, True),
], size=17)

# =============================================================== 12 in progress
s = prs.slides.add_slide(BLANK)
header(s, "In Progress — testing a semantics-preserving operator")
bullets(s, [
    ("Question: can the agent evade Safeline WITHOUT breaking the SQL?", 0, ACCENT, True),
    ("Added div_break (a /0 context-breaker) from a manual bypass PoC:", 0),
    ("-1' UNION SELECT database(),...  →  -1'/0 UNION SELECT database(),...   (extraction intact)", 1, MUT),
    ("Cheap PROBE first (no training): does a /0 payload BOTH pass Safeline AND return the marker?", 0),
    ("Retrain the agents only if the probe shows a reachable success.", 1, MUT),
], size=18)

# =============================================================== 13 next steps
s = prs.slides.add_slide(BLANK)
header(s, "Next Steps")
bullets(s, [
    ("1.  Run the /0 probe on Safeline → decide whether retraining is worthwhile.", 0),
    ("2.  Lock RQ2: three-regime result + the evade-vs-steal mechanism.", 0),
    ("3.  Writing (KIIT): add RQ1-error (mean±std) + the RQ2 Safeline section; update abstract.", 0, ACCENT, True),
    ("4.  Methods: document Safeline setup + WAF-evasion metric + limitations.", 0),
    ("Summary: RQ1 (union+error) ✓ · RQ3 ✓ · RQ2 transfer ✓ — stronger than planned.", 0, GREEN, True),
], size=19)

# =============================================================== Speaker notes
NOTES = [
 ("Slide pembuka. Sebut nama proyek, ini laporan progress, dan tanggal. Tujuan SeqSQLi: melatih agen "
  "RL untuk mem-bypass WAF lewat URUTAN mutasi.",
  "Opening. State the project, that this is a progress report, and the date. SeqSQLi trains an RL agent "
  "to bypass WAFs via a SEQUENCE of mutations."),
 ("Masalah inti: WAF memblok lewat pencocokan pola; payload bisa diubah permukaannya tanpa mengubah makna "
  "serangan. Kita modelkan sbg langkah-demi-langkah: tiap langkah pilih 1 mutasi; reward = bypass NYATA "
  "(lolos WAF + data marker keluar). 3 algoritma (PPO/TRPO/A2C), 2 WAF (ModSec rule & Safeline ML).",
  "Core problem: WAFs block by pattern-matching; a payload's surface can change without changing attack "
  "meaning. We model it step-by-step: each step picks one mutation; reward = a REAL bypass (passes WAF + "
  "extracts the marker). 3 algorithms, 2 WAFs (ModSec rule & Safeline ML)."),
 ("Tiga pertanyaan riset. RQ1 = algoritma mana terbaik. RQ2 = apakah policy dari WAF rule-based bisa "
  "dipindah ke WAF ML-based. RQ3 = apakah urutan mutasi penting. Semua sudah ada hasil; baru minggu ini "
  "= RQ1 error + RQ2 Safeline.",
  "Three research questions. RQ1 = which algorithm is best. RQ2 = does a rule-based policy move to an "
  "ML-based WAF. RQ3 = does mutation order matter. All have results; new this week = RQ1 error + RQ2 Safeline."),
 ("Slide referensi 2 metrik (jelaskan SEBELUM tabel hasil). IFNR = seberapa SERING agen tembus (tinggi = "
  "bagus); karena WAF blok semua payload mentah (FNR0=0), IFNR ≈ tingkat keberhasilan agen. SPBARC = "
  "seberapa HEMAT: rata-rata berapa kali coba (request) untuk 1 bypass (rendah = bagus). Agen ideal: "
  "IFNR tinggi + SPBARC rendah. TRPO menang dua-duanya.",
  "Metrics reference (explain BEFORE the result tables). IFNR = how OFTEN the agent breaks through "
  "(higher = better); since the WAF blocks all raw payloads (FNR0=0), IFNR ≈ the agent's success rate. "
  "SPBARC = how EFFICIENT: average tries (requests) per bypass (lower = better). Ideal agent: high IFNR "
  "+ low SPBARC. TRPO wins both."),
 ("Hasil utama RQ1 (union, ModSec). TRPO menang telak (IFNR 99%), lalu PPO, lalu A2C. Ketiganya kuasai "
  "tier mudah/medium; pembedanya tier COMPLEX (rantai mutasi terpanjang). SPBARC TRPO paling kecil = "
  "paling efisien.",
  "Main RQ1 result (union, ModSec). TRPO clearly wins (IFNR 99%), then PPO, then A2C. All master the "
  "easy/medium tiers; the differentiator is the COMPLEX tier (longest chains). TRPO's SPBARC is lowest = "
  "most efficient."),
 ("Hasil BARU: corpus error-based, tiap algoritma dites 3 kali (3 seed) biar adil. (1) Error jauh lebih "
  "sulit: IFNR terbaik anjlok 99%->30%. (2) Yang TIDAK konsisten di A2C = hasil/policy akhirnya: std ~14% "
  "(~7x TRPO/PPO), bisa 6% atau 31% cuma gara-gara seed acak. (3) Untungnya kita tes 3x: satu seed A2C "
  "sempat 0% — kalau dipakai sendiri, klaim 'A2C kolaps' itu SALAH.",
  "NEW: error-based corpus, each algorithm run 3 times (3 seeds) to be fair. (1) Error is much harder: best "
  "IFNR drops 99%->30%. (2) What is INCONSISTENT in A2C = its final policy/score: std ~14% (~7x TRPO/PPO), "
  "lands at 6% or 31% from the random seed alone. (3) Running 3 seeds mattered: one A2C seed was 0% — using "
  "it alone would have been a WRONG 'A2C collapses' claim."),
 ("Urutan mutasi penting. 35 pasangan consensus lintas-algoritma: membalik urutan menjatuhkan keberhasilan. "
  "Contoh precondition: ident_backtick HARUS sebelum hex_to_char. Maknanya: ordering = sinyal struktural/kausal "
  "dari topologi aturan WAF, bukan kebetulan satu algoritma.",
  "Mutation order matters. 35 cross-algorithm consensus pairs: reversing the order drops success. "
  "Precondition example: ident_backtick must precede hex_to_char. Meaning: ordering is a structural/causal "
  "signal from the WAF rule topology, not an algorithm artifact."),
 ("Zero-shot transfer = ambil model dilatih di ModSec, uji di Safeline TANPA latih ulang. Hasil 0% di "
  "SEMUA (12 model); ~900 payload termutasi semua kena 403. Penting: 0% ini ASLI, bukan bug — positive "
  "control buktikan Safeline bisa ditembus (PoC manual balik 200+data). Jadi 0% murni karena policy ModSec "
  "tak cocok ke WAF ML.",
  "Zero-shot transfer = take ModSec-trained models, test on Safeline WITHOUT retraining. 0% across ALL "
  "(12 models); ~900 mutated payloads all got 403. Important: this 0% is GENUINE, not a bug — a positive "
  "control shows Safeline IS bypassable (manual PoC returns 200+data). So 0% is purely a policy mismatch "
  "to the ML-WAF."),
 ("Slide TERPENTING. Tiga situasi: (a) ModSec regex: lolos DAN nyolong data (76-99%). (b) Safeline transfer: "
  "0/0. (c) Safeline dilatih langsung: lolos WAF ~89% TAPI nyolong data 0%. Inti: di WAF regex 'lolos' & "
  "'nyolong' menyatu; di WAF ML keduanya TERPISAH. 846/972 request = 200 tapi SQL_ERROR -> lolos hanya "
  "dengan merusak SQL.",
  "MOST IMPORTANT slide. Three settings: (a) ModSec regex: evade AND steal data (76-99%). (b) Safeline "
  "transfer: 0/0. (c) Safeline trained directly: passes WAF ~89% BUT steals data 0%. Core: on regex-WAFs "
  "evade & steal go together; on ML-WAFs they SPLIT. 846/972 requests = 200 but SQL_ERROR -> passing only "
  "by breaking the SQL."),
 ("Kenapa 'lolos tapi nihil'. Detektor ML menilai makna, bukan pola permukaan. Agar terlihat benign, agen "
  "konvergen ke satu rantai perusak: double-url-encode bikin UNION bukan keyword lagi; trik GBK merusak "
  "quote di backend UTF-8. Akibatnya Safeline maupun MySQL tak lihat SQL valid -> lolos (200) tapi error "
  "(tak ada data). Pelajaran: di WAF ML, agen hanya bisa lolos dengan mengorbankan serangannya sendiri.",
  "Why 'passes but nothing'. The ML detector judges meaning, not surface patterns. To look harmless, the "
  "agent converges to one destructive chain: double-URL-encoding makes UNION not a keyword; the GBK trick "
  "breaks the quote on a UTF-8 backend. So neither Safeline nor MySQL sees valid SQL -> passes (200) but "
  "errors (no data). Lesson: on ML-WAFs the agent can only evade by sacrificing its own attack."),
 ("Positioning JUJUR. BWAFSQLi (ACM TOSEM 2026) sudah mem-bypass Safeline dgn payload menjaga-semantik dan "
  "sudah mengalahkan metode RL. Jadi 'RL menemukan bypass Safeline' BUKAN kontribusi kita. Kebaruan kita: "
  "(1) ukur transferability gap rule->ML; (2) analisis ordering kausal; (3) decoupling lolos-vs-nyolong di "
  "WAF ML. Kita pelajari KENAPA policy struktural transfer/tidak, bukan bikin tool bypass baru.",
  "HONEST positioning. BWAFSQLi (ACM TOSEM 2026) already bypasses Safeline (semantic-preserving) and already "
  "beats RL methods. So 'RL finds a Safeline bypass' is NOT our contribution. Our novelty: (1) the rule->ML "
  "transferability gap; (2) causal ordering analysis; (3) evade-vs-steal decoupling on ML-WAFs. We study WHY "
  "structural policies transfer — not a new bypass tool."),
 ("Yang sedang dikerjakan. Pertanyaan: bisakah agen lolos Safeline TANPA merusak SQL? Kita tambah operator "
  "div_break (sisip /0) dari PoC manual, ekstraksi tetap utuh. Langkah hemat: jalankan PROBE dulu (tanpa "
  "training) untuk cek apakah payload /0 bisa lolos WAF sekaligus mengembalikan marker. Retrain hanya kalau "
  "probe menjanjikan.",
  "In progress. Question: can the agent evade Safeline WITHOUT breaking the SQL? We added div_break (insert "
  "/0) from the manual PoC, extraction intact. Cost-saving step: run a PROBE first (no training) to check "
  "whether a /0 payload passes the WAF AND returns the marker. Retrain only if the probe is promising."),
 ("Rencana: (1) jalankan probe /0 -> putuskan retrain. (2) Kunci RQ2 (3-rezim + mekanisme). (3) Menulis "
  "paper publikasi: tambah RQ1-error + section RQ2 Safeline, perbarui abstract. (4) Metodologi: dokumentasi "
  "setup Safeline + metrik evasi + limitasi. Ringkasan: RQ1(union+error), RQ3, RQ2 transfer SELESAI — lebih "
  "kuat dari rencana awal.",
  "Plan: (1) run the /0 probe -> decide on retraining. (2) Lock RQ2 (three-regime + mechanism). (3) Write "
  "the publication paper: add RQ1-error + the RQ2 Safeline section; update the abstract. (4) Methods: "
  "document Safeline setup + evasion metric + limitations. Summary: RQ1(union+error), RQ3, RQ2 transfer "
  "DONE — stronger than planned."),
]
for slide, (idn, eng) in zip(prs.slides, NOTES):
    tf = slide.notes_slide.notes_text_frame
    tf.text = "[ID] " + idn
    tf.add_paragraph().text = ""
    tf.add_paragraph().text = "[EN] " + eng

out = "SeqSQLi_progress_labmeeting.pptx"
try:
    prs.save(out)
except PermissionError:
    out = "SeqSQLi_progress_labmeeting_v2.pptx"
    print(f"[!] default file locked (open in PowerPoint?) -> saving to {out}", file=sys.stderr)
    prs.save(out)
print(f"[*] saved {out} with {len(prs.slides._sldIdLst)} slides")
